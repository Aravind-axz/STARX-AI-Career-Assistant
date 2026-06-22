# app.py — STARX AI Career Assistant
# Main Flask application entry point

import os
import json
import PyPDF2
import docx
from pptx import Presentation
from flask import Flask, render_template, request, jsonify, session
from werkzeug.utils import secure_filename

from notes_generator import generate_notes
from interview_simulator import get_interview_question, evaluate_answer, generate_final_report
from resume_analyzer import analyze_resume
from study_planner import generate_study_plan

# ── App Setup ───────────────────────────────────────────────────────────────
app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "starx-secret-dev-key-2024")

UPLOAD_FOLDER = "uploads"
ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx", "doc"}
MAX_CONTENT_LENGTH = 10 * 1024 * 1024  # 10 MB

app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["MAX_CONTENT_LENGTH"] = MAX_CONTENT_LENGTH

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# ── Helpers ──────────────────────────────────────────────────────────────────
def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_file(filepath: str) -> str:
    """Extracts plain text from PDF, DOCX, or PPTX files."""
    ext = filepath.rsplit(".", 1)[1].lower()
    text = ""

    if ext == "pdf":
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

    elif ext in ("docx", "doc"):
        doc = docx.Document(filepath)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

    elif ext == "pptx":
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"

    return text.strip()


def get_api_key() -> str:
    """Gets the Gemini API key from session or environment."""
    return session.get("api_key") or os.environ.get("GEMINI_API_KEY", "")


# ── Routes: Pages ────────────────────────────────────────────────────────────
@app.route("/")
def index():
    return render_template("index.html")

@app.route("/notes")
def notes():
    return render_template("notes.html")

@app.route("/interview")
def interview():
    return render_template("interview.html")

@app.route("/resume")
def resume():
    return render_template("resume.html")

@app.route("/planner")
def planner():
    return render_template("planner.html")


# ── API: API Key Management ───────────────────────────────────────────────────
@app.route("/api/set-key", methods=["POST"])
def set_api_key():
    """Stores the Gemini API key in the user session."""
    data = request.get_json()
    key = data.get("api_key", "").strip()
    if not key:
        return jsonify({"error": "API key is required"}), 400
    session["api_key"] = key
    return jsonify({"success": True, "message": "API key saved for this session"})


# ── API: Notes Generator ──────────────────────────────────────────────────────
@app.route("/api/notes/generate", methods=["POST"])
def api_generate_notes():
    """Accepts a document file, extracts text, and returns AI-generated notes."""
    api_key = get_api_key()
    if not api_key:
        return jsonify({"error": "Gemini API key not set. Please add it in settings."}), 401

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if not file or not allowed_file(file.filename):
        return jsonify({"error": "Invalid file type. Upload PDF, DOCX, or PPTX."}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    file.save(filepath)

    try:
        text = extract_text_from_file(filepath)
        if not text or len(text) < 100:
            return jsonify({"error": "Could not extract enough text from the file. Try a different document."}), 400

        notes = generate_notes(text, api_key)
        return jsonify({"success": True, "data": notes, "word_count": len(text.split())})

    except Exception as e:
        return jsonify({"error": f"AI generation failed: {str(e)}"}), 500

    finally:
        if os.path.exists(filepath):
            os.remove(filepath)


# ── API: Interview Simulator ──────────────────────────────────────────────────
@app.route("/api/interview/start", methods=["POST"])
def api_interview_start():
    """Starts a new interview session and returns the first question."""
    api_key = get_api_key()
    if not api_key:
        return jsonify({"error": "Gemini API key not set."}), 401

    data = request.get_json()
    role = data.get("role", "software_developer")

    # Initialize fresh session state
    session["interview_role"] = role
    session["interview_history"] = []
    session["interview_evaluations"] = []
    session["interview_question_num"] = 1

    try:
        result = get_interview_question(role, 1, [], api_key)
        session["interview_current_question"] = result["question"]
        return jsonify({"success": True, "data": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/interview/answer", methods=["POST"])
def api_interview_answer():
    """Receives an answer, evaluates it, and returns the next question."""
    api_key = get_api_key()
    if not api_key:
        return jsonify({"error": "Gemini API key not set."}), 401

    data = request.get_json()
    answer = data.get("answer", "").strip()

    if not answer:
        return jsonify({"error": "Answer cannot be empty"}), 400

    role = session.get("interview_role", "software_developer")
    history = session.get("interview_history", [])
    evaluations = session.get("interview_evaluations", [])
    question_num = session.get("interview_question_num", 1)
    current_question = session.get("interview_current_question", "")

    try:
        # Evaluate the submitted answer
        evaluation = evaluate_answer(role, current_question, answer, api_key)
        evaluations.append(evaluation)
        history.append({"question": current_question, "answer": answer})

        session["interview_evaluations"] = evaluations
        session["interview_history"] = history

        # Check if interview is complete
        if question_num >= 8:
            report = generate_final_report(role, evaluations, api_key)
            return jsonify({
                "success": True,
                "evaluation": evaluation,
                "interview_complete": True,
                "final_report": report
            })

        # Get next question
        next_q_num = question_num + 1
        session["interview_question_num"] = next_q_num
        next_q = get_interview_question(role, next_q_num, history, api_key)
        session["interview_current_question"] = next_q["question"]

        return jsonify({
            "success": True,
            "evaluation": evaluation,
            "interview_complete": False,
            "next_question": next_q
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── API: Resume Analyzer ──────────────────────────────────────────────────────
@app.route("/api/resume/analyze", methods=["POST"])
def api_resume_analyze():
    """Accepts a resume file and returns comprehensive AI analysis."""
    api_key = get_api_key()
    if not api_key:
        return jsonify({"error": "Gemini API key not set."}), 401

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if not file or not allowed_file(file.filename):
        return jsonify({"error": "Invalid file type. Upload PDF or DOCX."}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    file.save(filepath)

    try:
        text = extract_text_from_file(filepath)
        if not text or len(text) < 50:
            return jsonify({"error": "Could not extract text from resume. Please try a different format."}), 400

        result = analyze_resume(text, api_key)
        return jsonify({"success": True, "data": result})

    except Exception as e:
        return jsonify({"error": f"Analysis failed: {str(e)}"}), 500

    finally:
        if os.path.exists(filepath):
            os.remove(filepath)


# ── API: Study Planner ────────────────────────────────────────────────────────
@app.route("/api/planner/generate", methods=["POST"])
def api_planner_generate():
    """Generates a personalized study plan."""
    api_key = get_api_key()
    if not api_key:
        return jsonify({"error": "Gemini API key not set."}), 401

    data = request.get_json()
    exam_date = data.get("exam_date", "")
    subjects = data.get("subjects", [])
    daily_hours = int(data.get("daily_hours", 4))

    if not exam_date or not subjects:
        return jsonify({"error": "Exam date and subjects are required"}), 400

    try:
        result = generate_study_plan(exam_date, subjects, daily_hours, api_key)
        return jsonify({"success": True, "data": result})
    except Exception as e:
        return jsonify({"error": f"Plan generation failed: {str(e)}"}), 500


# ── Run ───────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    debug = os.environ.get("FLASK_ENV", "development") == "development"
    app.run(host="0.0.0.0", port=port, debug=debug)
